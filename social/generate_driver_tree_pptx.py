#!/usr/bin/env python3
"""
generate_driver_tree_pptx.py
Generates the Total Z-score Driver Tree as a Canva-importable .pptx file.
Every text box, rectangle and colour block lands as a separate editable element in Canva.

Usage:
    python3 generate_driver_tree_pptx.py
Output:
    social/z_score_driver_tree.pptx
"""

import os
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ── Palette ────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0a, 0x0e, 0x1a)
SURFACE = RGBColor(0x13, 0x18, 0x27)
DARK    = RGBColor(0x08, 0x0c, 0x16)
TEXT    = RGBColor(0xff, 0xff, 0xff)
MUTED   = RGBColor(0x88, 0x92, 0xa4)
GREEN   = RGBColor(0x00, 0xe6, 0x76)
BLUE    = RGBColor(0x4d, 0x9f, 0xff)
BORDER  = RGBColor(0x20, 0x27, 0x3d)
FACTOR  = RGBColor(0xcc, 0xd5, 0xe0)
CHEVRON = RGBColor(0x50, 0x60, 0x80)

# ── Data ───────────────────────────────────────────────────────────────────
STATS = [
    {'name': 'Points',     'role': ['Minutes', '2PA', '3PA', 'FTA'],  'rate': ['2P%', '3P%', 'FT%']},
    {'name': '3-Pointers', 'role': ['Minutes', '3PA'],                'rate': ['3P%']},
    {'name': 'Rebounds',   'role': ['Minutes'],                       'rate': ['Def. Reb Rate', 'Off. Reb Rate']},
    {'name': 'Assists',    'role': ['Minutes'],                       'rate': ['Assist Rate']},
    {'name': 'Steals',     'role': ['Minutes'],                       'rate': ['Steals Rate']},
    {'name': 'Blocks',     'role': ['Minutes'],                       'rate': ['Blocks Rate']},
    {'name': 'FG%',        'role': ['Minutes', '3PA', '2PA'],         'rate': ['2P%', '3P%']},
    {'name': 'FT%',        'role': ['Minutes', 'FTA'],                'rate': ['FT%']},
    {'name': 'Turnovers',  'role': ['Minutes'],                       'rate': ['Turnover Rate']},
]

# ── Layout (inches) ────────────────────────────────────────────────────────
W, H   = 10.0, 10.0     # slide size
LM     = 0.30            # left margin
RM     = 0.30            # right margin
TM     = 0.28            # top margin

N        = len(STATS)
COL_GAP  = 0.07
COL_W    = (W - LM - RM - (N - 1) * COL_GAP) / N   # ≈ 0.982"

GRID_Y       = 2.22
HEAD_H       = 0.64
ROLE_BAND_Y  = GRID_Y + HEAD_H           # 2.86
ROLE_BAND_H  = 0.44
ROLE_Y       = ROLE_BAND_Y + ROLE_BAND_H # 3.30
ROLE_H       = 2.10
RATE_BAND_Y  = ROLE_Y + ROLE_H           # 5.40
RATE_BAND_H  = 0.44
RATE_Y       = RATE_BAND_Y + RATE_BAND_H # 5.84
RATE_H       = 1.78

FONT = 'Calibri'


def col_x(i):
    return LM + i * (COL_W + COL_GAP)


# ── Low-level helpers ──────────────────────────────────────────────────────

def _no_line(shape):
    """Remove the outline from a shape via XML."""
    spPr = shape._element.spPr
    for ln in spPr.findall(qn('a:ln')):
        spPr.remove(ln)
    ln = etree.SubElement(spPr, qn('a:ln'))
    etree.SubElement(ln, qn('a:noFill'))


def add_rect(slide, x, y, w, h, fill, border_color=None, border_pt=0.6):
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border_color:
        shp.line.color.rgb = border_color
        shp.line.width = Pt(border_pt)
    else:
        _no_line(shp)
    return shp


def add_txt(slide, x, y, w, h, text, pt, color,
            bold=False, italic=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return tb


def add_mixed_txt(slide, x, y, w, h, parts, pt, align=PP_ALIGN.LEFT):
    """Single text box with multiple coloured runs — all in one paragraph."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = False
    para = tf.paragraphs[0]
    para.alignment = align
    for text, color, bold in parts:
        run = para.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(pt)
        run.font.color.rgb = color
        run.font.bold = bold
    return tb


def add_factors(slide, x, y, w, items, pt=8.5, line_h=0.22):
    """Stack labelled factor items with a chevron bullet."""
    for i, item in enumerate(items):
        iy = y + i * line_h
        add_txt(slide, x,        iy, 0.12, line_h, '›', pt + 1, CHEVRON)
        add_txt(slide, x + 0.12, iy, w - 0.12, line_h, item, pt, FACTOR)


# ── Build presentation ─────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)
slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

# Full background
bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(W), Inches(H))
bg.fill.solid()
bg.fill.fore_color.rgb = BG
_no_line(bg)

# ──────────────────────────────────────────────────────────────────────────
# TOP BAR
# ──────────────────────────────────────────────────────────────────────────

# Logo: emoji + Roto (white) + Intel (green) — single text box, 3 runs
add_mixed_txt(slide, LM, TM, 2.4, 0.32,
              [('🏀  ', TEXT, True), ('Roto', TEXT, True), (' Intel', GREEN, True)],
              pt=13.5)

# Badge pill
add_rect(slide, W - RM - 1.6, TM - 0.04, 1.6, 0.32, fill=BG, border_color=MUTED, border_pt=0.5)
add_txt(slide, W - RM - 1.6, TM + 0.02, 1.6, 0.22,
        'FANTASY ANALYTICS', 7, MUTED, bold=True, align=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────────────────────────────────
# TITLE BLOCK
# ──────────────────────────────────────────────────────────────────────────

add_txt(slide, 0, 0.75, W, 0.22,
        'HOW EACH CATEGORY IS EXPLAINED',
        8.5, GREEN, bold=True, align=PP_ALIGN.CENTER)

# Main title: "Total " + "Z-score" (green) + " Driver Tree"
add_mixed_txt(slide, 0, 1.00, W, 0.80,
              [('Total ', TEXT, True), ('Z-score', GREEN, True), (' Driver Tree', TEXT, True)],
              pt=38, align=PP_ALIGN.CENTER)

add_txt(slide, 0, 1.82, W, 0.26,
        'Every fantasy stat decomposed into the factors that drive it',
        9.5, MUTED, align=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────────────────────────────────
# COLUMN HEADERS (row 1)
# ──────────────────────────────────────────────────────────────────────────

for i, s in enumerate(STATS):
    x = col_x(i)
    add_rect(slide, x, GRID_Y, COL_W, HEAD_H, fill=SURFACE, border_color=BORDER)
    add_rect(slide, x, GRID_Y, COL_W, 0.026,  fill=GREEN)       # green top accent
    add_txt(slide, x, GRID_Y + 0.04, COL_W, HEAD_H - 0.04,
            s['name'], 10.5, TEXT, bold=True, align=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────────────────────────────────
# ROLE BAND LABEL
# ──────────────────────────────────────────────────────────────────────────

add_rect(slide, LM, ROLE_BAND_Y + 0.13, 0.26, 0.032, fill=BLUE)   # pip
add_txt(slide, LM + 0.34, ROLE_BAND_Y, W - LM - RM - 0.34, ROLE_BAND_H,
        'ROLE FACTORS  —  VOLUME & OPPORTUNITY', 8.5, BLUE, bold=True)

# ──────────────────────────────────────────────────────────────────────────
# ROLE CELLS (row 3)
# ──────────────────────────────────────────────────────────────────────────

for i, s in enumerate(STATS):
    x = col_x(i)
    add_rect(slide, x, ROLE_Y, COL_W, ROLE_H, fill=SURFACE, border_color=BORDER)
    add_rect(slide, x, ROLE_Y, COL_W, 0.024, fill=BLUE)           # blue top accent
    add_txt(slide, x + 0.07, ROLE_Y + 0.10, COL_W - 0.1, 0.18,
            'ROLE', 6, BLUE, bold=True)
    add_factors(slide, x + 0.07, ROLE_Y + 0.30, COL_W - 0.1, s['role'])

# ──────────────────────────────────────────────────────────────────────────
# RATE BAND LABEL
# ──────────────────────────────────────────────────────────────────────────

add_rect(slide, LM, RATE_BAND_Y + 0.13, 0.26, 0.032, fill=GREEN)  # pip
add_txt(slide, LM + 0.34, RATE_BAND_Y, W - LM - RM - 0.34, RATE_BAND_H,
        'RATE FACTORS  —  EFFICIENCY & SKILL', 8.5, GREEN, bold=True)

# ──────────────────────────────────────────────────────────────────────────
# RATE CELLS (row 5)
# ──────────────────────────────────────────────────────────────────────────

for i, s in enumerate(STATS):
    x = col_x(i)
    add_rect(slide, x, RATE_Y, COL_W, RATE_H, fill=DARK, border_color=BORDER)
    add_rect(slide, x, RATE_Y, COL_W, 0.024, fill=GREEN)          # green top accent
    add_txt(slide, x + 0.07, RATE_Y + 0.10, COL_W - 0.1, 0.18,
            'RATE', 6, GREEN, bold=True)
    add_factors(slide, x + 0.07, RATE_Y + 0.30, COL_W - 0.1, s['rate'])

# ──────────────────────────────────────────────────────────────────────────
# SAVE
# ──────────────────────────────────────────────────────────────────────────

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'z_score_driver_tree.pptx')
prs.save(out)
print(f'✓  Saved → {out}')
